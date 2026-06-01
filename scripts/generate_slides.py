#!/usr/bin/env python3
"""Generate GWAS PrePubMatch hackathon slide deck."""

from pathlib import Path
from typing import List, Optional, Tuple

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

REPO_ROOT = Path(__file__).resolve().parent.parent
OUTPUT = REPO_ROOT / "slides" / "GWAS_PrePubMatch.pptx"

# Brand palette (matches AGENTS.md diagram colors)
INDIGO = RGBColor(0x63, 0x66, 0xF1)
EMERALD = RGBColor(0x10, 0xB9, 0x81)
SKY = RGBColor(0x0E, 0xA5, 0xE9)
DARK = RGBColor(0x1E, 0x29, 0x3B)
SLATE = RGBColor(0x64, 0x74, 0x8B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF8, 0xFA, 0xFC)


def _set_slide_bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_title_bar(slide, title: str, subtitle: Optional[str] = None) -> None:
    bar = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        Inches(0),
        Inches(0),
        Inches(13.333),
        Inches(1.35),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = INDIGO
    bar.line.fill.background()
    tf = bar.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT
    tf.margin_left = Inches(0.6)
    tf.margin_top = Inches(0.25)
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(16)
        p2.font.color.rgb = RGBColor(0xE0, 0xE7, 0xFF)
        p2.space_before = Pt(4)


def _add_bullets(slide, items: List[str], top: float = 1.7, left: float = 0.7) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(11.8), Inches(5.5))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(22)
        p.font.color.rgb = DARK
        p.space_after = Pt(14)
        p.bullet = True


def _add_flow_boxes(slide, boxes: List[Tuple[str, RGBColor]], top: float = 2.0) -> None:
    n = len(boxes)
    width = 2.0
    gap = 0.35
    total = n * width + (n - 1) * gap
    start = (13.333 - total) / 2
    for i, (label, color) in enumerate(boxes):
        left = start + i * (width + gap)
        shape = slide.shapes.add_shape(
            1,
            Inches(left),
            Inches(top),
            Inches(width),
            Inches(1.1),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        tf = shape.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        if i < n - 1:
            arrow = slide.shapes.add_textbox(
                Inches(left + width),
                Inches(top + 0.35),
                Inches(gap),
                Inches(0.5),
            )
            ap = arrow.text_frame.paragraphs[0]
            ap.text = "→"
            ap.font.size = Pt(20)
            ap.font.color.rgb = SLATE
            ap.alignment = PP_ALIGN.CENTER


def slide_title(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _set_slide_bg(slide, LIGHT_BG)

    accent = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(0.12))
    accent.fill.solid()
    accent.fill.fore_color.rgb = INDIGO
    accent.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.5), Inches(1.5))
    tp = title_box.text_frame.paragraphs[0]
    tp.text = "GWAS PrePubMatch"
    tp.font.size = Pt(48)
    tp.font.bold = True
    tp.font.color.rgb = DARK

    sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(3.2), Inches(11.5), Inches(1.2))
    sp = sub_box.text_frame.paragraphs[0]
    sp.text = "Discover pre-publication GWAS summary statistics\nalongside published Catalog studies and literature"
    sp.font.size = Pt(24)
    sp.font.color.rgb = SLATE

    tag_box = slide.shapes.add_textbox(Inches(0.8), Inches(5.5), Inches(11.5), Inches(0.6))
    tagp = tag_box.text_frame.paragraphs[0]
    tagp.text = "AGI House Hackathon · May 2026"
    tagp.font.size = Pt(18)
    tagp.font.color.rgb = EMERALD


def slide_problem(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    _add_title_bar(slide, "The Problem")
    _add_bullets(
        slide,
        [
            "~8,000 pre-publication GWAS summary statistics in the GWAS Catalog are poorly indexed for search",
            "Researchers submit messy metadata — draft titles, author lists, free-text trait strings",
            "No unified view linking pre-pub sumstats to published studies and their literature",
            "Manual reconciliation is slow, error-prone, and blocks reproducible discovery",
        ],
    )


def slide_solution(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    _add_title_bar(slide, "Our Solution")
    _add_bullets(
        slide,
        [
            "Paste submission metadata → get ranked Catalog studies + publications in one view",
            "Deterministic match scoring (no invented PMIDs, DOIs, or GCST accessions)",
            "Parallel search: GWAS Catalog Solr + OpenAlex, Europe PMC, PubMed",
            "Production-oriented FastAPI orchestrator + React dashboard",
        ],
    )


def slide_architecture(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    _add_title_bar(slide, "Architecture", "FastAPI orchestrator — not a simple proxy")

    _add_flow_boxes(
        slide,
        [
            ("React UI", INDIGO),
            ("FastAPI", INDIGO),
            ("GWAS Catalog", EMERALD),
            ("Literature", SKY),
            ("Rank & Merge", RGBColor(0x06, 0xB6, 0xD4)),
        ],
        top=1.9,
    )

    _add_flow_boxes(
        slide,
        [
            ("Science Skills", RGBColor(0x8B, 0x5C, 0xF6)),
            ("OpenAlex", SKY),
            ("Europe PMC", SKY),
            ("PubMed", SKY),
        ],
        top=3.5,
    )

    note = slide.shapes.add_textbox(Inches(0.7), Inches(5.0), Inches(11.8), Inches(1.5))
    np = note.text_frame.paragraphs[0]
    np.text = (
        "Literature dual-path: Google Science Skills CLIs (primary) with direct HTTP fallback per source. "
        "Returns unified JSON (schema v3.1.0) with full provenance."
    )
    np.font.size = Pt(18)
    np.font.color.rgb = SLATE


def slide_features(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    _add_title_bar(slide, "Key Features")
    _add_bullets(
        slide,
        [
            "Partial failure tolerance — Catalog results even if literature sources degrade",
            "Health probes — /api/health checks each upstream with latency",
            "TTL response cache — configurable for Solr and literature",
            "Transparent provenance — match scores, source status, skills vs HTTP backend",
            "Docker deployment — Science Skills vendored at build time",
        ],
    )


def slide_demo(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    _add_title_bar(slide, "Demo Flow")
    _add_bullets(
        slide,
        [
            "Load a scenario preset (Probable Match, High Confidence, Ambiguous Trait, No Match)",
            "Edit title, authors, reported trait, GCST/DOI clues",
            "Click Discover Related Studies → unified ranked hits",
            "Review top-match verification links (PubMed, GWAS Catalog, DOI)",
            "Supplementary pre-pub sumstats panel for related unpublished entries",
        ],
        top=1.7,
    )

    ex = slide.shapes.add_textbox(Inches(0.7), Inches(5.3), Inches(11.8), Inches(1.2))
    ep = ex.text_frame.paragraphs[0]
    ep.text = 'Example: "Shared and distinct genetic risk factors for childhood-onset asthma" → Pividori asthma match'
    ep.font.size = Pt(16)
    ep.font.italic = True
    ep.font.color.rgb = EMERALD


def slide_closing(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, INDIGO)

    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11.5), Inches(1.0))
    tp = title_box.text_frame.paragraphs[0]
    tp.text = "Try It"
    tp.font.size = Pt(44)
    tp.font.bold = True
    tp.font.color.rgb = WHITE

    _add_bullets(
        slide,
        [
            "Local: bash scripts/run_server.sh  +  cd web && npm run dev",
            "Docker: docker compose up --build  →  http://localhost:8080",
            "API: POST /api/discover  ·  GET /api/health",
        ],
        top=3.5,
        left=0.8,
    )
    for shape in slide.shapes:
        if shape.has_text_frame and shape != title_box:
            for p in shape.text_frame.paragraphs:
                p.font.color.rgb = WHITE


def main() -> None:
    OUTPUT.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_title(prs)
    slide_problem(prs)
    slide_solution(prs)
    slide_architecture(prs)
    slide_features(prs)
    slide_demo(prs)
    slide_closing(prs)

    prs.save(str(OUTPUT))
    print(f"Wrote {OUTPUT} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
